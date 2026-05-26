# ingest.py ───────────────────────────────────────────────────────────────────
# Generic ingestion pipeline — works for ANY document type:
#   product catalogs, SOPs, HR policies, training manuals, etc.
#
# Key design principle:
#   Metadata is DERIVED from the document itself (headings, filename, structure).
#   Nothing is hardcoded about what topics or products exist.
# ─────────────────────────────────────────────────────────────────────────────

import os
import re
import time
import hashlib
import requests
from collections import defaultdict
from dotenv import load_dotenv

from docx import Document
import pdfplumber
from pptx import Presentation
from pinecone import Pinecone

load_dotenv()

# =============================================================================
# CONFIG
# =============================================================================

DOCS_FOLDER  = r"C:\Training GPT\Document"   # 👈 folder with your files

HF_API_KEY   = os.getenv("HF_API_KEY")
HF_API_URL   = "https://router.huggingface.co/hf-inference/models/BAAI/bge-base-en-v1.5"
HF_HEADERS   = {
    "Authorization": f"Bearer {HF_API_KEY}",
    "Content-Type":  "application/json",
}

PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
PINECONE_INDEX   = "sales-chatbot"

CHUNK_SIZE    = 800
CHUNK_OVERLAP = 200

print("HF_API_KEY    :", "✅" if HF_API_KEY       else "❌ MISSING")
print("PINECONE_KEY  :", "✅" if PINECONE_API_KEY  else "❌ MISSING")


# =============================================================================
# METADATA HELPERS
# Derived from document content — no hardcoded domain knowledge.
# =============================================================================

def infer_doc_category(filename: str) -> str:
    """
    Infer a broad document category from the filename.
    This is a loose label derived purely from the filename — not a fixed taxonomy.
    Add more patterns here as your document library grows.
    """
    name = filename.lower()
    if any(k in name for k in ["sop", "procedure", "process"]):
        return "sop"
    if any(k in name for k in ["policy", "policies", "hr", "leave", "conduct", "compliance"]):
        return "policy"
    if any(k in name for k in ["product", "catalog", "sofa", "mattress", "chair", "recliner"]):
        return "product"
    if any(k in name for k in ["training", "onboarding", "manual", "guide", "handbook"]):
        return "training"
    if any(k in name for k in ["price", "pricing", "rate", "tariff"]):
        return "pricing"
    if any(k in name for k in ["faq", "question", "answer"]):
        return "faq"
    return "general"   # safe fallback — never crashes


def make_chunk_id(source: str, heading: str, chunk_index: int) -> str:
    """Deterministic ID — re-running ingest never creates duplicates."""
    raw = f"{source}::{heading}::{chunk_index}"
    return hashlib.md5(raw.encode()).hexdigest()


# =============================================================================
# EXTRACTORS
# Each returns a list of section dicts:
#   { heading, text, source, file_type, doc_category }
#
# Headings come directly from the document structure — not from any predefined
# list. Whatever heading the document author wrote becomes the section label.
# =============================================================================

def extract_docx(filepath: str) -> list[dict]:
    """
    Splits on Heading-styled paragraphs.
    Tables are flattened to 'Label: Value' lines for clean embedding.
    The heading text from the document itself becomes the section label.
    """
    doc          = Document(filepath)
    source       = os.path.basename(filepath)
    doc_category = infer_doc_category(os.path.basename(filepath))
    sections     = []

    # Use filename (without extension) as the default heading for content
    # before the first document heading
    current_heading = os.path.splitext(os.path.basename(source))[0]
    current_lines   = []

    def flush():
        text = "\n".join(current_lines).strip()
        if text:
            sections.append({
                "heading":      current_heading,
                "text":         text,
                "source":       source,
                "file_type":    "docx",
                "doc_category": doc_category,
            })

    from docx.oxml.ns import qn

    for block in doc.element.body:
        tag = block.tag.split("}")[-1]

        if tag == "p":
            style_elem = block.find(qn("w:pPr"))
            style_name = ""
            if style_elem is not None:
                s = style_elem.find(qn("w:pStyle"))
                if s is not None:
                    style_name = s.get(qn("w:val"), "")

            text = "".join(
                n.text for n in block.iter()
                if n.tag.split("}")[-1] == "t" and n.text
            ).strip()

            if not text:
                continue

            if style_name.lower().startswith("heading"):
                flush()
                current_heading = text
                current_lines   = []
            else:
                current_lines.append(text)

        elif tag == "tbl":
            # Flatten table rows as "Col1: Col2: Col3" — readable by the LLM
            for row in block.findall(
                ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tr"
            ):
                cells = row.findall(
                    ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tc"
                )
                cell_texts = []
                for cell in cells:
                    ct = "".join(
                        n.text for n in cell.iter()
                        if n.tag.split("}")[-1] == "t" and n.text
                    ).strip()
                    if ct:
                        cell_texts.append(ct)
                if cell_texts:
                    current_lines.append(": ".join(cell_texts))

    flush()
    return sections


def extract_pdf(filepath: str) -> list[dict]:
    """
    Each page becomes one section. Heading = 'Page N'.
    For PDFs with clear visual section headers, consider upgrading to
    pdfplumber's bbox analysis to detect headings from font size.
    """
    source       = os.path.basename(filepath)
    doc_category = infer_doc_category(source)
    sections     = []

    with pdfplumber.open(filepath) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text()
            if text and text.strip():
                sections.append({
                    "heading":      f"Page {i}",
                    "text":         text.strip(),
                    "source":       source,
                    "file_type":    "pdf",
                    "doc_category": doc_category,
                })
    return sections


def extract_pptx(filepath: str) -> list[dict]:
    """
    Each slide becomes one section.
    Heading = the slide's first non-empty text block (usually the slide title).
    """
    source       = os.path.basename(filepath)
    doc_category = infer_doc_category(source)
    prs          = Presentation(filepath)
    sections     = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts   = []
        heading = f"Slide {slide_num}"

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if not line:
                    continue
                if heading == f"Slide {slide_num}":
                    heading = line   # first text block becomes the heading
                else:
                    parts.append(line)

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Notes]: {notes}")

        if parts:
            sections.append({
                "heading":      heading,
                "text":         "\n".join(parts),
                "source":       source,
                "file_type":    "pptx",
                "doc_category": doc_category,
            })

    return sections


EXTRACTORS = {
    ".docx": extract_docx,
    ".pdf":  extract_pdf,
    ".pptx": extract_pptx,
}


# =============================================================================
# FILE LOADER
# =============================================================================

def load_all_files(folder_path: str) -> list[dict]:
    all_sections = []
    skipped      = []

    for filename in sorted(os.listdir(folder_path)):
        ext = os.path.splitext(filename)[1].lower()
        if ext not in EXTRACTORS:
            skipped.append(filename)
            continue

        filepath = os.path.join(folder_path, filename)
        try:
            print(f"  📄 [{ext[1:].upper()}]  {filename} ...", end=" ")
            sections = EXTRACTORS[ext](filepath)
            if not sections:
                print("⚠️  Empty — skipped")
                continue
            all_sections.extend(sections)
            char_total = sum(len(s["text"]) for s in sections)
            cat        = sections[0]["doc_category"]
            print(f"✅  {len(sections)} sections · {char_total:,} chars · [{cat}]")
        except Exception as e:
            print(f"❌  ERROR — {e}")

    if skipped:
        print(f"\n  ⏭️  Skipped (unsupported): {', '.join(skipped)}")

    return all_sections


# =============================================================================
# CHUNKER
# Sections that fit within CHUNK_SIZE are kept whole.
# Oversized sections are split with overlap so context is never lost at seams.
# =============================================================================

def split_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    chunks = []
    start  = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        if end >= len(text):
            break
        start += chunk_size - overlap
    return chunks


def chunk_all_sections(sections: list[dict]) -> tuple[list[str], list[dict]]:
    """
    Produces flat (text, metadata) pairs ready for embedding.

    Metadata fields stored in Pinecone per chunk:
      id            — deterministic content hash (safe to re-run)
      source        — original filename
      heading       — section heading taken from the document itself
      doc_category  — loose type inferred from filename
      file_type     — docx / pdf / pptx
      chunk_index   — position within the section
      text          — the chunk content (also stored for retrieval)
    """
    texts     = []
    metadatas = []
    summary   = defaultdict(int)

    for section in sections:
        heading      = section["heading"]
        body         = section["text"]
        source       = section["source"]
        file_type    = section["file_type"]
        doc_category = section["doc_category"]

        sub_chunks = split_text(body) if len(body) > CHUNK_SIZE else [body]

        for idx, chunk in enumerate(sub_chunks):
            chunk_id = make_chunk_id(source, heading, idx)
            texts.append(chunk)
            metadatas.append({
                "id":           chunk_id,
                "source":       source,
                "heading":      heading,
                "doc_category": doc_category,
                "file_type":    file_type,
                "chunk_index":  idx,
                "text":         chunk,
            })

        summary[doc_category] += len(sub_chunks)
        label = f"[{doc_category:12s}] {source[:25]:25s} / {heading[:35]}"
        print(f"  ✂️  {label}  → {len(sub_chunks)} chunk(s)")

    print("\n  📊 Chunk summary by doc_category:")
    for cat, count in summary.items():
        print(f"     {cat.upper():15s}: {count} chunks")

    return texts, metadatas


# =============================================================================
# EMBEDDING
# Context prefix uses doc_category + heading from the document itself.
# MUST stay in sync with build_query_embed_text() in app.py.
#
# Format:
#   "Category: <doc_category>. Heading: <heading>.\n<chunk text>"
# =============================================================================

def build_embed_text(chunk: str, meta: dict) -> str:
    return (
        f"Category: {meta['doc_category']}. "
        f"Heading: {meta['heading']}.\n"
        f"{chunk}"
    )


def get_embedding(text: str) -> list[float]:
    response = requests.post(HF_API_URL, headers=HF_HEADERS, json={"inputs": text})
    result   = response.json()

    if isinstance(result, dict):
        if "error" in result:
            raise Exception(f"HF API Error: {result['error']}")
        if "estimated_time" in result:
            raise Exception("model_loading")

    if isinstance(result, list):
        if isinstance(result[0], float):        return result
        if isinstance(result[0], list):
            if isinstance(result[0][0], float): return result[0]
            if isinstance(result[0][0], list):  return result[0][0]

    raise ValueError(f"Unexpected HF response: {str(result)[:200]}")


def get_embedding_with_retry(text: str, retries: int = 5) -> list[float]:
    wait = 5
    for attempt in range(1, retries + 1):
        try:
            return get_embedding(text)
        except Exception as e:
            msg = str(e).lower()
            if "model_loading" in msg or "503" in msg or "loading" in msg:
                print(f"  ⏳ HF model loading — retrying in {wait}s ({attempt}/{retries})")
                time.sleep(wait)
                wait = min(wait * 2, 60)
            else:
                raise
    raise Exception(f"Embedding failed after {retries} retries")


def embed_all(texts: list[str], metadatas: list[dict]) -> list[list[float]]:
    embeddings = []
    total      = len(texts)
    for i, (chunk, meta) in enumerate(zip(texts, metadatas)):
        embed_text = build_embed_text(chunk, meta)
        emb        = get_embedding_with_retry(embed_text)
        embeddings.append(emb)
        if (i + 1) % 10 == 0 or (i + 1) == total:
            print(f"  🔢 Embedded {i + 1} / {total}")
    return embeddings


# =============================================================================
# PINECONE UPLOAD
# =============================================================================

def upload_to_pinecone(embeddings: list[list[float]], metadatas: list[dict]) -> None:
    pc    = Pinecone(api_key=PINECONE_API_KEY)
    index = pc.Index(PINECONE_INDEX)
    print(f"  ✅ Connected → {PINECONE_INDEX}")

    records = [
        {
            "id":     meta["id"],
            "values": emb,
            "metadata": {
                "text":         meta["text"],
                "source":       meta["source"],
                "heading":      meta["heading"],
                "doc_category": meta["doc_category"],
                "file_type":    meta["file_type"],
                "chunk_index":  meta["chunk_index"],
            },
        }
        for emb, meta in zip(embeddings, metadatas)
    ]

    batch_size = 100
    for i in range(0, len(records), batch_size):
        index.upsert(vectors=records[i : i + batch_size])
        print(f"  ☁️  Upserted {min(i + batch_size, len(records))} / {len(records)}")

    print(f"\n  ✅ Upload complete — {len(records)} vectors")


# =============================================================================
# MAIN
# =============================================================================

if __name__ == "__main__":
    print("=" * 60)
    print("  GENERIC RAG — INGESTION PIPELINE")
    print("=" * 60)

    if not os.path.isdir(DOCS_FOLDER):
        print(f"\n❌  Folder not found: {DOCS_FOLDER}")
        exit(1)

    print(f"\n📂  Scanning: {DOCS_FOLDER}")
    sections = load_all_files(DOCS_FOLDER)
    print(f"\n  → {len(sections)} sections extracted")

    if not sections:
        print("\n❌  No content found. Check folder path and file formats.")
        exit(1)

    print(f"\n✂️   Chunking ...")
    texts, metadatas = chunk_all_sections(sections)
    print(f"\n  → {len(texts)} total chunks")

    print(f"\n🧠  Embedding via HF API ...")
    embeddings = embed_all(texts, metadatas)

    print(f"\n☁️   Uploading to Pinecone ...")
    upload_to_pinecone(embeddings, metadatas)

    categories_seen = set(m["doc_category"] for m in metadatas)
    sources_seen    = set(m["source"] for m in metadatas)

    print("\n" + "=" * 60)
    print("  ✅  INGESTION COMPLETE")
    print("=" * 60)
    print(f"  Files      : {len(sources_seen)}")
    print(f"  Chunks     : {len(texts)}")
    print(f"  Categories : {', '.join(sorted(categories_seen))}")
    print(f"  Sources    : {', '.join(sorted(sources_seen))}")
    print("=" * 60)